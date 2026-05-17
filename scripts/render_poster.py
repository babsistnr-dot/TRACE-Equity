"""Rendert bericht/poster.html via Headless-Chrome zu einem
vektoriellen A0-PDF (das Poster ist für Druck gebaut: @media print,
@page A0, .poster 841×1189 mm) und rastert daraus eine
hochauflösende PNG, die als vollflächiges Bild auf genau einer
A0-Folie ins PPTX eingebettet wird.

WICHTIG: Der Screenshot-Weg funktioniert NICHT — das Poster nutzt
on-screen ein (in der sauberen poster.html entferntes) Skalier-JS;
nur der Druckpfad (--print-to-pdf) löst die korrekte A0-Print-CSS aus.

    python scripts/render_poster.py            # PDF + PNG + PPTX
    python scripts/render_poster.py --pdf-only # nur A0-PDF

Voraussetzungen: Google Chrome; PyMuPDF (PDF→PNG); python-pptx.
Fehlende Pakete werden gemeldet.
"""

import os
import subprocess
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
POSTER = ROOT / "Poster"
POSTER_HTML = POSTER / "poster.html"
PDF_OUT = POSTER / "poster_render.pdf"
PNG_OUT = POSTER / "poster_render.png"
PPTX_OUT = POSTER / "Poster_TRACE-Equity.pptx"

# A0 Hochformat
A0_W_MM, A0_H_MM = 841, 1189
RASTER_DPI = 300                        # PDF→PNG für PPTX-Einbettung

CHROME_CANDIDATES = [
    r"C:\Program Files\Google\Chrome\Application\chrome.exe",
    r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
    os.path.expandvars(r"%LOCALAPPDATA%\Google\Chrome\Application\chrome.exe"),
    r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
    r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
]


def find_chrome():
    for c in CHROME_CANDIDATES:
        if Path(c).is_file():
            return c
    raise SystemExit("Chrome/Edge nicht gefunden.")


def render_pdf():
    if not POSTER_HTML.is_file():
        raise SystemExit(f"Fehlt: {POSTER_HTML} — erst extract_poster_html.py.")
    chrome = find_chrome()
    uri = POSTER_HTML.resolve().as_uri()
    args = [
        chrome,
        "--headless=new",
        "--disable-gpu",
        "--no-sandbox",
        "--no-pdf-header-footer",
        "--virtual-time-budget=15000",
        f"--print-to-pdf={PDF_OUT}",
        uri,
    ]
    subprocess.run(args, check=False, capture_output=True)
    if not PDF_OUT.is_file():
        raise SystemExit("Render fehlgeschlagen (kein PDF erzeugt).")
    print(f"[OK] {PDF_OUT.relative_to(ROOT)} (vektoriell, A0 via Print-CSS)")


def rasterize():
    try:
        import fitz  # PyMuPDF
    except ModuleNotFoundError:
        raise SystemExit("PyMuPDF fehlt → 'pip install pymupdf', dann erneut.")
    doc = fitz.open(PDF_OUT)
    page = doc[0]
    zoom = RASTER_DPI / 72.0
    pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
    pix.save(PNG_OUT)
    dpi = round(pix.width / (A0_W_MM / 25.4))
    print(f"[OK] {PNG_OUT.relative_to(ROOT)}  {pix.width}×{pix.height}px "
          f"(~{dpi} dpi @ A0)")
    if dpi < 150:
        print("  [WARN] < 150 dpi — RASTER_DPI erhöhen.")


def build_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Mm
    except ModuleNotFoundError:
        raise SystemExit("python-pptx fehlt → 'pip install python-pptx', "
                          "dann erneut (oder --png-only).")
    prs = Presentation()
    prs.slide_width = Mm(A0_W_MM)
    prs.slide_height = Mm(A0_H_MM)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # leeres Layout
    slide.shapes.add_picture(str(PNG_OUT), 0, 0,
                             width=Mm(A0_W_MM), height=Mm(A0_H_MM))
    prs.save(PPTX_OUT)
    print(f"[OK] {PPTX_OUT.relative_to(ROOT)}  (1 Folie, A0, Bild vollflächig)")


if __name__ == "__main__":
    render_pdf()
    if "--pdf-only" not in sys.argv:
        rasterize()
        build_pptx()
